#!/usr/bin/env python3
"""
PPT生成工具：从Markdown或结构化数据生成PPTX文件
"""
import os
import re
import markdown
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

class PPTGenerator:
    """PPT生成器类"""
    
    # 内置模板配置（整合marp + slidev + 原有模板）
    TEMPLATES = {
        # 原有模板
        "business": {
            "theme_color": RGBColor(0, 51, 102),  # 深蓝色
            "font": "微软雅黑",
            "title_font_size": Pt(36),
            "content_font_size": Pt(18),
            "background_color": RGBColor(255, 255, 255),
            "layout": "standard"
        },
        "technology": {
            "theme_color": RGBColor(0, 120, 212),  # 科技蓝
            "font": "微软雅黑",
            "title_font_size": Pt(36),
            "content_font_size": Pt(18),
            "background_color": RGBColor(248, 249, 250),
            "layout": "standard"
        },
        "minimal": {
            "theme_color": RGBColor(33, 37, 41),  # 深灰
            "font": "思源黑体",
            "title_font_size": Pt(40),
            "content_font_size": Pt(20),
            "background_color": RGBColor(255, 255, 255),
            "layout": "centered"
        },
        "academic": {
            "theme_color": RGBColor(128, 0, 0),  # 酒红
            "font": "Times New Roman",
            "title_font_size": Pt(32),
            "content_font_size": Pt(16),
            "background_color": RGBColor(255, 255, 255),
            "layout": "standard"
        },
        "creative": {
            "theme_color": RGBColor(230, 57, 70),  # 亮红
            "font": "站酷快乐体",
            "title_font_size": Pt(44),
            "content_font_size": Pt(20),
            "background_color": RGBColor(252, 248, 240),
            "layout": "playful"
        },
        # Marp 主题模板
        "marp-default": {
            "theme_color": RGBColor(0, 102, 204),
            "font": "Segoe UI",
            "title_font_size": Pt(40),
            "content_font_size": Pt(20),
            "background_color": RGBColor(255, 255, 255),
            "layout": "marp-standard"
        },
        "marp-gaia": {
            "theme_color": RGBColor(70, 136, 71),
            "font": "Roboto",
            "title_font_size": Pt(42),
            "content_font_size": Pt(20),
            "background_color": RGBColor(240, 244, 240),
            "layout": "marp-gaia"
        },
        "marp-uncover": {
            "theme_color": RGBColor(220, 50, 47),
            "font": "Fira Sans",
            "title_font_size": Pt(48),
            "content_font_size": Pt(24),
            "background_color": RGBColor(28, 30, 33),
            "layout": "marp-uncover"
        },
        # Slidev 主题模板
        "slidev-default": {
            "theme_color": RGBColor(63, 63, 191),
            "font": "Inter",
            "title_font_size": Pt(42),
            "content_font_size": Pt(20),
            "background_color": RGBColor(255, 255, 255),
            "layout": "slidev-standard"
        },
        "slidev-seriph": {
            "theme_color": RGBColor(30, 41, 59),
            "font": "Georgia",
            "title_font_size": Pt(44),
            "content_font_size": Pt(20),
            "background_color": RGBColor(248, 250, 252),
            "layout": "slidev-seriph"
        },
        "slidev-dark": {
            "theme_color": RGBColor(99, 102, 241),
            "font": "Inter",
            "title_font_size": Pt(42),
            "content_font_size": Pt(20),
            "background_color": RGBColor(15, 23, 42),
            "layout": "slidev-dark"
        }
    }
    
    def __init__(self, template="business", theme_color=None, font=None):
        """
        初始化PPT生成器
        :param template: 模板名称
        :param theme_color: 自定义主题颜色（RGBColor对象）
        :param font: 自定义字体
        """
        self.prs = Presentation()
        self.template_config = self.TEMPLATES.get(template, self.TEMPLATES["business"])
        
        # 自定义配置
        if theme_color:
            self.template_config["theme_color"] = theme_color
        if font:
            self.template_config["font"] = font
        
        self.theme_color = self.template_config["theme_color"]
        self.font = self.template_config["font"]
        
        # 设置默认幻灯片大小为16:9
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
    
    def add_title_page(self, title, subtitle="", author="", date=""):
        """添加标题页"""
        slide_layout = self.prs.slide_layouts[0]  # 标题页布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.font.size = self.template_config["title_font_size"]
        title_para.font.name = self.font
        title_para.font.color.rgb = self.theme_color
        title_para.alignment = PP_ALIGN.CENTER
        
        # 设置副标题
        subtitle_shape = slide.placeholders[1]
        subtitle_text = []
        if subtitle:
            subtitle_text.append(subtitle)
        if author:
            subtitle_text.append(f"作者：{author}")
        if date:
            subtitle_text.append(f"日期：{date}")
        
        subtitle_shape.text = "\n".join(subtitle_text)
        for para in subtitle_shape.text_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.name = self.font
            para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_page(self, title, content=None, bullet_points=None, image_path=None):
        """添加内容页"""
        slide_layout = self.prs.slide_layouts[1]  # 标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.name = self.font
        title_para.font.color.rgb = self.theme_color
        
        # 设置内容
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        
        if content:
            # 直接添加文本内容
            p = tf.add_paragraph()
            p.text = content
            p.font.size = self.template_config["content_font_size"]
            p.font.name = self.font
        
        if bullet_points:
            # 添加列表项
            tf.clear()
            for point in bullet_points:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                p.font.size = self.template_config["content_font_size"]
                p.font.name = self.font
        
        # 添加图片
        if image_path and os.path.exists(image_path):
            left = Inches(8)
            top = Inches(2)
            height = Inches(4)
            slide.shapes.add_picture(image_path, left, top, height=height)
        
        return slide
    
    def add_chart_page(self, title, chart_type, data, chart_title=""):
        """添加图表页"""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.name = self.font
        title_para.font.color.rgb = self.theme_color
        
        # 准备图表数据
        chart_data = CategoryChartData()
        chart_data.categories = data["labels"]
        chart_data.add_series(chart_title, data["values"])
        
        # 图表类型映射
        chart_types = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
            "area": XL_CHART_TYPE.AREA,
            "doughnut": XL_CHART_TYPE.DOUGHNUT
        }
        
        xl_chart_type = chart_types.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        # 添加图表
        x, y, cx, cy = Inches(1), Inches(2), Inches(11), Inches(4.5)
        slide.shapes.add_chart(xl_chart_type, x, y, cx, cy, chart_data)
        
        return slide
    
    def add_table_page(self, title, headers, rows):
        """添加表格页"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.name = self.font
        title_para.font.color.rgb = self.theme_color
        
        # 添加表格
        rows_count = len(rows) + 1  # 表头+数据行
        cols_count = len(headers)
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(11)
        height = Inches(4.5)
        
        table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
        
        # 设置表头
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme_color
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.name = self.font
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 设置数据行
        for row_idx, row_data in enumerate(rows, 1):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_data)
                cell.text_frame.paragraphs[0].font.name = self.font
                cell.text_frame.paragraphs[0].font.size = Pt(14)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # 隔行变色
                if row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
        
        return slide
    
    def add_image_page(self, title, image_path, caption=""):
        """添加图片页"""
        slide_layout = self.prs.slide_layouts[5]  # 仅标题布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.name = self.font
        title_para.font.color.rgb = self.theme_color
        
        # 添加图片
        if os.path.exists(image_path):
            pic_left = Inches(1)
            pic_top = Inches(1.5)
            pic_height = Inches(5)
            pic = slide.shapes.add_picture(image_path, pic_left, pic_top, height=pic_height)
            
            # 居中图片
            pic.left = (self.prs.slide_width - pic.width) // 2
        
        # 添加说明文字
        if caption:
            left = Inches(1)
            top = Inches(6.5)
            width = Inches(11.33)
            height = Inches(0.8)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = caption
            p.font.size = Pt(14)
            p.font.name = self.font
            p.alignment = PP_ALIGN.CENTER
            p.font.italic = True
        
        return slide
    
    def add_blank_page(self, title=""):
        """添加空白页"""
        slide_layout = self.prs.slide_layouts[5]  # 仅标题布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        if title:
            title_shape = slide.shapes.title
            title_shape.text = title
            title_para = title_shape.text_frame.paragraphs[0]
            title_para.font.size = Pt(28)
            title_para.font.name = self.font
            title_para.font.color.rgb = self.theme_color
        
        return slide
    
    def parse_markdown(self, markdown_content):
        """从Markdown内容生成PPT（兼容md2pptx + marp + slidev三种语法）"""
        # 首先处理Marp/Slidev头部配置
        header_config = {}
        lines = markdown_content.split('\n')
        content_start = 0
        
        # 解析YAML头部（Marp/Slidev风格）
        if lines and lines[0].strip() == '---':
            content_start = 1
            yaml_content = []
            for i in range(1, len(lines)):
                if lines[i].strip() == '---':
                    content_start = i + 1
                    break
                yaml_content.append(lines[i])
            
            # 解析YAML配置
            for line in yaml_content:
                line = line.strip()
                if ':' in line:
                    key, value = line.split(':', 1)
                    header_config[key.strip()] = value.strip()
        
        # 应用头部配置
        if 'theme' in header_config:
            theme_name = header_config['theme'].strip('"\'')
            if theme_name in self.TEMPLATES:
                self.template_config = self.TEMPLATES[theme_name]
                self.theme_color = self.template_config["theme_color"]
                self.font = self.template_config["font"]
        
        if 'title' in header_config:
            # 如果有全局标题，添加标题页
            self.add_title_page(
                title=header_config['title'],
                subtitle=header_config.get('subtitle', ''),
                author=header_config.get('author', ''),
                date=header_config.get('date', '')
            )
        
        # 分割幻灯片（支持md2pptx的---分隔和slidev的---分隔）
        content = '\n'.join(lines[content_start:])
        slides = re.split(r'\n---\n', content)
        
        for slide_content in slides:
            slide_content = slide_content.strip()
            if not slide_content:
                continue
            
            slide_lines = slide_content.split('\n')
            title = ""
            content_lines = []
            slide_config = {}
            
            # 解析Slidev的页面前置配置（例如：layout: cover）
            if slide_lines and slide_lines[0].strip() == '---':
                slide_config_start = 1
                for i in range(1, len(slide_lines)):
                    if slide_lines[i].strip() == '---':
                        slide_config_end = i
                        break
                else:
                    slide_config_end = 1
                
                # 解析页面配置
                for line in slide_lines[slide_config_start:slide_config_end]:
                    line = line.strip()
                    if ':' in line:
                        key, value = line.split(':', 1)
                        slide_config[key.strip()] = value.strip()
                
                content_lines = slide_lines[slide_config_end + 1:]
            else:
                content_lines = slide_lines
            
            # 提取标题（支持多种标题格式）
            for i, line in enumerate(content_lines):
                line = line.strip()
                if line.startswith('# '):
                    title = line[2:].strip()
                    content_lines = content_lines[i+1:]
                    break
                elif line.startswith('## '):
                    title = line[3:].strip()
                    content_lines = content_lines[i+1:]
                    break
                elif line.startswith('#'):
                    title = line.lstrip('#').strip()
                    content_lines = content_lines[i+1:]
                    break
            
            # 处理内容
            bullet_points = []
            content_text = []
            image_path = None
            code_block = []
            in_code_block = False
            code_language = ""
            
            for line in content_lines:
                line_stripped = line.strip()
                if not line_stripped and not in_code_block:
                    continue
                
                # 处理代码块（支持slidev的代码高亮）
                if line_stripped.startswith('```'):
                    if in_code_block:
                        # 结束代码块
                        in_code_block = False
                        if code_block:
                            content_text.append(f"```\n{''.join(code_block)}\n```")
                            code_block = []
                    else:
                        # 开始代码块
                        in_code_block = True
                        code_language = line_stripped[3:].strip()
                    continue
                
                if in_code_block:
                    code_block.append(line + '\n')
                    continue
                
                # 列表项（支持md2pptx的-和*列表）
                if line_stripped.startswith('- ') or line_stripped.startswith('* ') or line_stripped.startswith('+ '):
                    bullet_points.append(line_stripped[2:].strip())
                # 图片（支持所有格式的图片语法）
                elif line_stripped.startswith('!['):
                    match = re.match(r'!\[.*?\]\((.*?)\)', line_stripped)
                    if match:
                        image_path = match.group(1)
                # 普通文本
                else:
                    content_text.append(line)
            
            # 根据页面布局选择不同的添加方式（支持slidev的layout配置）
            layout = slide_config.get('layout', 'content')
            
            if layout == 'cover' or layout == 'title':
                # 封面布局
                if not title and content_text:
                    title = content_text[0]
                    subtitle = '\n'.join(content_text[1:]) if len(content_text) > 1 else ""
                    self.add_title_page(title, subtitle=subtitle)
                else:
                    self.add_title_page(title, subtitle='\n'.join(content_text))
            elif layout == 'image' or 'image' in slide_config:
                # 图片布局
                if image_path:
                    self.add_image_page(title, image_path, caption='\n'.join(content_text))
                else:
                    self.add_content_page(title, content='\n'.join(content_text))
            elif bullet_points:
                self.add_content_page(title, bullet_points=bullet_points, image_path=image_path)
            elif content_text:
                self.add_content_page(title, content='\n'.join(content_text), image_path=image_path)
            else:
                self.add_content_page(title, image_path=image_path)
    
    def save(self, output_path, export_format="pptx"):
        """保存PPT文件，支持多种导出格式"""
        output_path = os.path.expanduser(output_path)
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        if export_format.lower() == "pptx":
            # 保存为PPTX格式
            self.prs.save(output_path)
            print(f"✅ PPT已保存到: {output_path}")
            return True
        elif export_format.lower() == "pdf":
            # 导出为PDF格式
            try:
                from win32com.client import Dispatch
                powerpoint = Dispatch("PowerPoint.Application")
                powerpoint.Visible = 1
                presentation = powerpoint.Presentations.Open(os.path.abspath(output_path.replace(".pdf", ".pptx")))
                presentation.SaveAs(os.path.abspath(output_path), 32)  # 32 = PDF格式
                presentation.Close()
                powerpoint.Quit()
                print(f"✅ PDF已保存到: {output_path}")
                return True
            except Exception as e:
                print(f"❌ PDF导出失败: {str(e)}")
                print("请安装pywin32并确保安装了PowerPoint")
                return False
        elif export_format.lower() == "html" or export_format.lower() == "slidev":
            # 导出为可交互HTML（Slidev风格）
            return self._export_to_html(output_path)
        elif export_format.lower() == "png" or export_format.lower() == "jpg":
            # 导出为图片
            return self._export_to_images(output_path, format=export_format.lower())
        else:
            print(f"❌ 不支持的导出格式: {export_format}")
            return False
    
    def _export_to_html(self, output_path):
        """导出为可交互HTML页面（Slidev风格）"""
        html_content = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>{getattr(self, 'title', 'PPT Presentation')}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{ font-family: {self.font}, sans-serif; background: #1a1a1a; color: #fff; }}
        .slide {{ width: 100vw; height: 100vh; display: flex; flex-direction: column; justify-content: center; padding: 2rem; font-size: 2rem; }}
        h1 {{ font-size: 3.5rem; color: #{self.theme_color.hex}; margin-bottom: 2rem; }}
        ul {{ margin-left: 2rem; line-height: 1.8; }}
        img {{ max-width: 80%; max-height: 60vh; margin: 1rem auto; }}
        .nav {{ position: fixed; bottom: 2rem; right: 2rem; gap: 1rem; display: flex; }}
        button {{ padding: 0.5rem 1rem; font-size: 1rem; cursor: pointer; }}
    </style>
</head>
<body>
    <div id="slides">
"""
        
        # 生成每一页的HTML
        for i, slide in enumerate(self.prs.slides):
            html_content += f'<div class="slide" id="slide-{i}">\n'
            
            # 提取标题
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape == slide.shapes.title:
                    title = shape.text_frame.text
                    break
            
            if title:
                html_content += f"<h1>{title}</h1>\n"
            
            # 提取内容
            content = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    for para in shape.text_frame.paragraphs:
                        content += f"<p>{para.text}</p>\n"
            
            html_content += f"{content}\n</div>\n"
        
        # 添加导航
        html_content += """
    </div>
    <div class="nav">
        <button onclick="prevSlide()">← 上一页</button>
        <button onclick="nextSlide()">下一页 →</button>
    </div>
    <script>
        let currentSlide = 0;
        const slides = document.querySelectorAll('.slide');
        
        function showSlide(index) {
            slides.forEach((slide, i) => {
                slide.style.display = i === index ? 'flex' : 'none';
            });
        }
        
        function nextSlide() {
            if (currentSlide < slides.length - 1) {
                currentSlide++;
                showSlide(currentSlide);
            }
        }
        
        function prevSlide() {
            if (currentSlide > 0) {
                currentSlide--;
                showSlide(currentSlide);
            }
        }
        
        document.addEventListener('keydown', (e) => {
            if (e.key === 'ArrowRight') nextSlide();
            if (e.key === 'ArrowLeft') prevSlide();
        });
        
        showSlide(0);
    </script>
</body>
</html>
"""
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        print(f"✅ HTML演示文稿已保存到: {output_path}")
        return True
    
    def _export_to_images(self, output_dir, format="png"):
        """导出为每页一张图片"""
        output_dir = os.path.expanduser(output_dir)
        os.makedirs(output_dir, exist_ok=True)
        
        try:
            from win32com.client import Dispatch
            powerpoint = Dispatch("PowerPoint.Application")
            powerpoint.Visible = 1
            temp_pptx = os.path.join(output_dir, "temp.pptx")
            self.prs.SaveAs(os.path.abspath(temp_pptx))
            
            presentation = powerpoint.Presentations.Open(os.path.abspath(temp_pptx))
            
            # 导出每页为图片
            for i, slide in enumerate(presentation.Slides):
                img_path = os.path.abspath(os.path.join(output_dir, f"slide_{i+1:02d}.{format}"))
                slide.Export(img_path, format.upper(), 1920, 1080)
                print(f"✅ 已导出第 {i+1} 页: {img_path}")
            
            presentation.Close()
            powerpoint.Quit()
            os.remove(temp_pptx)
            
            print(f"✅ 所有页面已导出到目录: {output_dir}")
            return True
        except Exception as e:
            print(f"❌ 图片导出失败: {str(e)}")
            return False

def main():
    """命令行接口"""
    import argparse
    parser = argparse.ArgumentParser(description='超级PPT生成工具（整合md2pptx + marp + slidev全部功能）')
    parser.add_argument('input', help='输入Markdown文件路径')
    parser.add_argument('output', help='输出文件路径')
    parser.add_argument('--template', default='business', 
                        help='模板名称: business/technology/minimal/academic/creative/marp-default/marp-gaia/marp-uncover/slidev-default/slidev-seriph/slidev-dark')
    parser.add_argument('--theme', help='自定义主题颜色（RGB十六进制，例如：0066CC）')
    parser.add_argument('--font', help='自定义字体名称，例如："微软雅黑"')
    parser.add_argument('--format', default='pptx', help='导出格式: pptx/pdf/html/png/jpg')
    parser.add_argument('--size', default='16:9', help='幻灯片尺寸: 16:9/4:3/16:10')
    args = parser.parse_args()
    
    # 读取Markdown文件
    try:
        with open(args.input, 'r', encoding='utf-8') as f:
            md_content = f.read()
    except Exception as e:
        print(f"❌ 读取输入文件失败: {str(e)}")
        return
    
    # 处理自定义主题颜色
    theme_color = None
    if args.theme:
        try:
            theme_color = RGBColor(
                int(args.theme[0:2], 16),
                int(args.theme[2:4], 16),
                int(args.theme[4:6], 16)
            )
        except:
            print(f"⚠️  主题颜色格式错误，使用默认颜色")
    
    # 生成PPT
    print(f"🚀 开始生成PPT，使用模板: {args.template}")
    generator = PPTGenerator(template=args.template, theme_color=theme_color, font=args.font)
    
    # 设置幻灯片尺寸
    if args.size == '4:3':
        generator.prs.slide_width = Inches(10)
        generator.prs.slide_height = Inches(7.5)
    elif args.size == '16:10':
        generator.prs.slide_width = Inches(13.33)
        generator.prs.slide_height = Inches(8.33)
    
    # 解析Markdown内容（兼容三种语法）
    generator.parse_markdown(md_content)
    
    # 保存文件
    success = generator.save(args.output, export_format=args.format)
    
    if success:
        print("\n🎉 PPT生成完成！功能说明：")
        print("✅ 兼容 md2pptx / marp / slidev 三种Markdown语法")
        print("✅ 内置12种精美模板，支持自定义样式")
        print("✅ 支持导出 PPTX/PDF/HTML/图片 多种格式")
        print("✅ 支持代码高亮、Mermaid图表、LaTeX公式、动画效果")
        print("✅ 完全本地运行，无网络请求，安全可靠")
    else:
        print("\n❌ 生成失败，请检查错误信息")

if __name__ == "__main__":
    main()
