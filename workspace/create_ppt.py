# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色
PRIMARY = RGBColor(255, 152, 0)      # 橙色
SECONDARY = RGBColor(255, 193, 7)    # 黄色
ACCENT = RGBColor(76, 175, 80)       # 绿色
TEXT = RGBColor(55, 71, 79)          # 深灰
WHITE = RGBColor(255, 255, 255)

def add_title_slide(title, subtitle=""):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()
    
    # 装饰圆
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(4), Inches(4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = SECONDARY
    circle.line.fill.background()
    
    # 主标题
    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(32)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, items, emoji=""):
    """内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # 顶部条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    
    # 标题背景
    title_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.4), Inches(12.3), Inches(1.1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = SECONDARY
    title_bg.line.fill.background()
    
    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(11.9), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}" if emoji else title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.2))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        if len(item) < 50:
            p.font.size = Pt(28)
        else:
            p.font.size = Pt(24)
        p.font.color.rgb = TEXT
        p.space_after = Pt(16)

def add_story_slide(title, story_text, page_num=""):
    """故事页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 浅黄背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 248, 225)
    bg.line.fill.background()
    
    # 云朵装饰
    cloud = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(10.5), Inches(0.3), Inches(2.2), Inches(1.2))
    cloud.fill.solid()
    cloud.fill.fore_color.rgb = WHITE
    cloud.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(9), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"📖 {title}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # 故事卡片
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(12), Inches(5.3))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = PRIMARY
    card.line.width = Pt(3)
    
    # 故事文字
    story_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.2), Inches(4.8))
    tf = story_box.text_frame
    tf.word_wrap = True
    
    lines = story_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        if len(line) < 20:
            p.font.size = Pt(26)
        else:
            p.font.size = Pt(22)
        p.font.color.rgb = TEXT
        if '小兔子' in line or '太阳' in line:
            p.font.bold = True
            p.font.color.rgb = PRIMARY
        p.space_after = Pt(12)
    
    # 页码
    if page_num:
        page_box = slide.shapes.add_textbox(Inches(12), Inches(6.9), Inches(1), Inches(0.4))
        tf = page_box.text_frame
        p = tf.paragraphs[0]
        p.text = page_num
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT

def add_vocab_slide(title, vocab_list):
    """生字学习页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 白色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"✨ {title}"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # 生字卡片
    card_width = Inches(3.8)
    card_height = Inches(2.2)
    start_x = Inches(0.6)
    start_y = Inches(1.5)
    
    card_colors = [
        (255, 205, 210), (255, 224, 178), (200, 230, 201),
        (187, 222, 251), (225, 190, 231), (255, 236, 179)
    ]
    
    for i, vocab in enumerate(vocab_list):
        row = i // 3
        col = i % 3
        x = start_x + col * (card_width + Inches(0.25))
        y = start_y + row * (card_height + Inches(0.3))
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*card_colors[i % len(card_colors)])
        card.line.fill.background()
        
        # 汉字
        char_box = slide.shapes.add_textbox(x, y + Inches(0.1), card_width, Inches(1))
        tf = char_box.text_frame
        p = tf.paragraphs[0]
        p.text = vocab['char']
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = TEXT
        p.alignment = PP_ALIGN.CENTER
        
        # 拼音
        pinyin_box = slide.shapes.add_textbox(x, y + Inches(1.1), card_width, Inches(0.5))
        tf = pinyin_box.text_frame
        p = tf.paragraphs[0]
        p.text = vocab['pinyin']
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT
        p.alignment = PP_ALIGN.CENTER
        
        # 意思
        mean_box = slide.shapes.add_textbox(x, y + Inches(1.5), card_width, Inches(0.6))
        tf = mean_box.text_frame
        p = tf.paragraphs[0]
        p.text = vocab['meaning']
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT
        p.alignment = PP_ALIGN.CENTER

def add_question_slide():
    """互动问答页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 243, 224)
    bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🤔 动脑筋想一想"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # 问题卡片
    questions = [
        {"q": "小兔子什么时候去上学？", "a": "清晨", "emoji": "🌅"},
        {"q": "小兔子上学的心情怎么样？", "a": "高高兴兴", "emoji": "😊"},
        {"q": "太阳公公怎么了？", "a": "笑了", "emoji": "☀️"}
    ]
    
    y_pos = Inches(1.5)
    for i, q in enumerate(questions):
        # 问题框
        q_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                       Inches(0.8), y_pos, Inches(11.7), Inches(1.6))
        q_box.fill.solid()
        q_box.fill.fore_color.rgb = WHITE
        q_box.line.color.rgb = PRIMARY
        q_box.line.width = Pt(2)
        
        # 问题文字
        q_text = slide.shapes.add_textbox(Inches(1), y_pos + Inches(0.2), Inches(8), Inches(0.6))
        tf = q_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i+1}. {q['q']}"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT
        
        # 答案
        a_text = slide.shapes.add_textbox(Inches(1), y_pos + Inches(0.85), Inches(8), Inches(0.5))
        tf = a_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"   答案：{q['emoji']} {q['a']}"
        p.font.size = Pt(22)
        p.font.color.rgb = PRIMARY
        
        y_pos += Inches(1.8)

# ============================================
# 创建PPT页面
# ============================================

# 第1页：封面
add_title_slide("📚 快乐阅读小课堂", "第一课：小兔子上学记")

# 第2页：学习目标
add_content_slide("今天我们要学", [
    "🐰 认识新朋友 —— 小兔子",
    "📖 一起读故事，一起学本领",
    "✨ 认识新的字词朋友",
    "🎤 有感情地朗读故事"
], "🎯")

# 第3页：故事内容
add_story_slide("故事开始啦", 
    "清晨，太阳公公笑了。\n小兔子背上小书包，\n高高兴兴去上学。\n\n\"老师早！同学早！\"\n小兔子笑眯眯地说。",
    "1/2")

# 第4页：生字学习
add_vocab_slide("认识新字词", [
    {"char": "太", "pinyin": "tài", "meaning": "太阳 - 给我们光和热"},
    {"char": "阳", "pinyin": "yáng", "meaning": "太阳 - 白天在天上"},
    {"char": "笑", "pinyin": "xiào", "meaning": "笑了 - 开心、高兴"},
    {"char": "学", "pinyin": "xué", "meaning": "上学 - 去学校学习"},
    {"char": "师", "pinyin": "shī", "meaning": "老师 - 教我们知识"},
    {"char": "早", "pinyin": "zǎo", "meaning": "早上 - 一天的开头"}
])

# 第5页：朗读指导
add_content_slide("一起读一读", [
    "\"清晨，太阳公公笑了。\"",
    "读的时候要注意：",
    "",
    "✅ 声音响亮，让全班都听见",
    "✅ 有感情，像讲故事一样",
    "✅ 不添字、不漏字、不错字"
], "🎤")

# 第6页：互动问答
add_question_slide()

# 第7页：收获总结
add_content_slide("今天的收获", [
    "1️⃣ 读了一个有趣的故事《小兔子上学记》",
    "",
    "2️⃣ 认识了6个新的字词朋友",
    "",
    "3️⃣ 学会了有感情地朗读",
    "",
    "4️⃣ 明白了上学要开开心心的道理"
], "🌟")

# 第8页：作业布置
add_content_slide("作业小任务", [
    "📚 把今天的故事读给爸爸妈妈听",
    "",
    "🖍️ 用彩笔画一幅\"小兔子上学\"的图画",
    "",
    "💬 和爸爸妈妈聊聊：你上学时是什么心情？"
], "📝")

# 第9页：结束页
add_title_slide("下节课再见！", "👋 谢谢小朋友们！")

# 保存PPT
output_path = r"C:\Users\LWS\.openclaw\workspace\一年级语文阅读第一课.pptx"
prs.save(output_path)
print(f"PPT generated successfully: {output_path}")
