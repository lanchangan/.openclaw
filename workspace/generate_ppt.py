#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成一年级阅读理解第一课PPT
适合6-7岁小学生，色彩丰富、字体大、内容简洁
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿（16:9）
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色方案（温暖明亮的颜色，适合小学生）
COLORS = {
    'coral': RGBColor(255, 107, 107),      # 珊瑚红
    'teal': RGBColor(78, 205, 196),        # 青绿色
    'yellow': RGBColor(255, 217, 61),      # 明黄
    'purple': RGBColor(155, 89, 182),      # 紫色
    'orange': RGBColor(230, 126, 34),      # 橙色
    'green': RGBColor(46, 204, 113),       # 绿色
    'blue': RGBColor(52, 152, 219),        # 蓝色
    'pink': RGBColor(255, 105, 180),       # 粉红
    'dark': RGBColor(44, 62, 80),          # 深色
    'white': RGBColor(255, 255, 255),      # 白色
}

def add_title_slide(prs, title, subtitle):
    """创建标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景装饰圆形
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(4), Inches(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['coral']
    shape.line.fill.background()
    
    shape2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), Inches(5), Inches(3), Inches(3))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = COLORS['teal']
    shape2.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS['coral']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(36)
    p.font.color.rgb = COLORS['teal']
    p.alignment = PP_ALIGN.CENTER
    
    # Emoji装饰
    emoji_box = slide.shapes.add_textbox(Inches(5.5), Inches(5.5), Inches(2.333), Inches(0.8))
    tf = emoji_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📚 ✨ 📖"
    p.font.size = Pt(48)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, color_key='blue', emoji="💡"):
    """创建内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 左侧色块装饰
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLORS[color_key]
    left_bar.line.fill.background()
    
    # 顶部标题背景
    title_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.3), Inches(12), Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLORS[color_key]
    title_bg.line.fill.background()
    
    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(32)
        p.font.color.rgb = COLORS['dark']
        p.space_before = Pt(12)
        p.line_spacing = 1.3
    
    return slide

def add_end_slide(prs, title, message):
    """创建结束页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景装饰
    shape1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(5), Inches(5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = COLORS['yellow']
    shape1.fill.transparency = 0.3
    shape1.line.fill.background()
    
    shape2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(4), Inches(4))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = COLORS['teal']
    shape2.fill.transparency = 0.3
    shape2.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = COLORS['coral']
    p.alignment = PP_ALIGN.CENTER
    
    # 信息
    msg_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(1.5))
    tf = msg_box.text_frame
    p = tf.paragraphs[0]
    p.text = message
    p.font.size = Pt(40)
    p.font.color.rgb = COLORS['teal']
    p.alignment = PP_ALIGN.CENTER
    
    # 装饰emoji
    emoji_box = slide.shapes.add_textbox(Inches(5), Inches(5.2), Inches(3.333), Inches(1))
    tf = emoji_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🌟 📚 🏆 🎉"
    p.font.size = Pt(52)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ===================== 生成PPT =====================

# 第1页：封面
add_title_slide(prs, "📚 阅读理解第一课", "🌟 我是小小阅读家 🌟")

# 第2页：学习目标
add_content_slide(prs, "今天我们要学习", [
    "👀 学会看图片 —— 仔细观察图画中的内容",
    "👂 学会听故事 —— 认真听老师讲故事",
    "🗣️ 学会讲故事 —— 用自己的话讲出故事",
    "✏️ 学会做练习 —— 回答简单的问题"
], "blue", "🎯")

# 第3页：认识封面
add_content_slide(prs, "第一步：认识封面", [
    "📖 看书的时候，先从封面开始看！",
    "",
    "🔍 封面上有什么？",
    "    📌 书名（这本书叫什么名字）",
    "    🎨 图画（画了什么内容）",
    "    ✍️ 作者（谁写了这本书）"
], "purple", "📕")

# 第4页：看图说话
add_content_slide(prs, "第二步：看图说话", [
    "🖼️ 看图有方法！记住这个顺序：",
    "",
    "    👉 从左到右",
    "    ⬇️ 从上到下",
    "",
    "🤔 看图时想一想：",
    "    ❓ 图上有什么？（人物、动物、景物）",
    "    ❓ 他们在做什么？（动作、表情）",
    "    ❓ 你觉得会发生什么？（想象）"
], "teal", "👀")

# 第5页：听故事
add_content_slide(prs, "第三步：听故事", [
    "👂 认真听故事，记住这5个W：",
    "",
    "    👤 Who  谁 —— 故事的主角是谁？",
    "    📍 Where 哪里 —— 故事发生在哪里？",
    "    ⏰ When  什么时候 —— 故事发生在什么时间？",
    "    ❓ What  什么 —— 故事里发生了什么事？",
    "    ✅ How   怎么样 —— 故事的结局是什么？"
], "orange", "📚")

# 第6页：回答问题
add_content_slide(prs, "第四步：回答问题", [
    "✏️ 做阅读理解练习的步骤：",
    "",
    "    1️⃣ 先看问题 —— 知道要问什么",
    "    2️⃣ 回书中找 —— 带着问题看书",
    "    3️⃣ 圈出答案 —— 找到关键词",
    "    4️⃣ 完整回答 —— 写清楚、写完整",
    "",
    "💡 小秘诀：答案往往在书里，仔细找！"
], "green", "📝")

# 第7页：分享故事
add_content_slide(prs, "第五步：分享故事", [
    "🗣️ 把故事讲给小伙伴听：",
    "",
    "    🎯 开头这样说：",
    "       \"今天我要讲一个关于...的故事\"",
    "",
    "    📖 中间讲清楚：",
    "       \"故事里发生了...（事情经过）\"",
    "",
    "    💡 结尾谈感受：",
    "       \"这个故事让我知道了...（道理/感受）\""
], "purple", "🌈")

# 第8页：阅读小技巧
add_content_slide(prs, "阅读小技巧总结", [
    "📚 记住这些阅读小法宝：",
    "",
    "    🔍 看封面，找信息",
    "    👀 看图画，有顺序",
    "    👂 听故事，抓重点",
    "    ✏️ 做练习，找答案",
    "    🗣️ 讲故事，表清楚",
    "",
    "🌟 每天读一点，你会越来越棒！"
], "pink", "⭐")

# 第9页：结束页
add_end_slide(prs, "🎉 你真棒！", "继续加油，成为最棒的小小阅读家！")

# 保存PPT
output_path = r"C:\Users\LWS\.openclaw\workspace\一年级阅读理解第一课.pptx"
prs.save(output_path)
print(f"PPT已成功生成！")
print(f"保存路径: {output_path}")
